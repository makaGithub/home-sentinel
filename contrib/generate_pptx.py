#!/usr/bin/env python3
"""
Генерация презентации проекта home-sentinel из presentation.md

Использование:
    python3 contrib/generate_pptx.py

Скрипт читает docs/presentation.md и генерирует docs/presentation.pptx
"""

import os
import re
import sys
import argparse

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.dml.fill import FillFormat
    print("✅ Все необходимые модули импортированы успешно")
except ImportError as e:
    print(f"❌ Ошибка импорта модулей: {e}")
    print("   Установите python-pptx: pip install python-pptx")
    sys.exit(1)

# Импортируем вспомогательный модуль для конфигурации
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
try:
    from config_helper import (
        load_env_file, get_project_root, generate_output_filename,
        get_build_dir, publish_to_gh_pages, create_github_release
    )
except ImportError as e:
    print(f"⚠️  Предупреждение: не удалось импортировать config_helper: {e}")
    print("   Используются значения по умолчанию")

# Константы форматирования
FONT_NAME = 'Arial'  # Arial лучше работает в презентациях
TITLE_SIZE = 44
SUBTITLE_SIZE = 32
HEADING_SIZE = 28
BODY_SIZE = 18
CODE_SIZE = 14

# Цвета - яркая палитра с градиентами
TITLE_COLOR = RGBColor(102, 126, 234)  # Яркий фиолетово-синий
TEXT_COLOR = RGBColor(51, 51, 51)  # Темно-серый для читаемости
ACCENT_COLOR = RGBColor(118, 75, 162)  # Фиолетовый
SUBTITLE_COLOR = RGBColor(102, 126, 234)  # Синий
CODE_BG_COLOR = RGBColor(245, 245, 250)  # Светло-серый для кода
BULLET_COLOR = RGBColor(102, 126, 234)  # Синий для маркеров

# Градиентные цвета для фона
GRADIENT_START = RGBColor(102, 126, 234)  # Синий
GRADIENT_END = RGBColor(118, 75, 162)  # Фиолетовый

def parse_markdown(md_file):
    """Парсит markdown файл и возвращает список слайдов"""
    with open(md_file, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # Разделяем на слайды по ---
    slides_content = content.split('---')
    
    slides = []
    for slide_content in slides_content:
        slide_content = slide_content.strip()
        if not slide_content:
            continue
        
        lines = slide_content.split('\n')
        slide_data = {
            'title': None,
            'subtitle': None,
            'content': [],
            'is_code': False
        }
        
        current_block = []
        in_code_block = False
        
        for line in lines:
            line = line.strip()
            
            if not line:
                if current_block:
                    slide_data['content'].append(('\n'.join(current_block), in_code_block))
                    current_block = []
                continue
            
            # Код блоки (```)
            if line.startswith('```'):
                if current_block:
                    slide_data['content'].append(('\n'.join(current_block), in_code_block))
                    current_block = []
                in_code_block = not in_code_block
                continue
            
            # Заголовки
            if line.startswith('# '):
                if current_block:
                    slide_data['content'].append(('\n'.join(current_block), in_code_block))
                    current_block = []
                slide_data['title'] = line[2:].strip()
            elif line.startswith('## '):
                if current_block:
                    slide_data['content'].append(('\n'.join(current_block), in_code_block))
                    current_block = []
                slide_data['subtitle'] = line[3:].strip()
            else:
                current_block.append(line)
        
        if current_block:
            slide_data['content'].append(('\n'.join(current_block), in_code_block))
        
        if slide_data['title'] or slide_data['content']:
            slides.append(slide_data)
    
    return slides

def calculate_font_size(text, max_width, max_height, initial_size, is_code=False):
    """Вычисляет оптимальный размер шрифта для текста"""
    # Подсчитываем количество строк
    lines = text.split('\n')
    lines_count = len([l for l in lines if l.strip()])  # Только непустые строки
    
    # Оцениваем количество символов
    total_chars = len(text)
    
    # Примерная оценка: для Arial 18pt примерно 8-10 символов на дюйм ширины
    # Межстрочный интервал примерно 1.2-1.5 от размера шрифта
    chars_per_line_estimate = max(1, int(max_width / (initial_size * 0.01)))  # Примерно
    estimated_lines = max(lines_count, total_chars / max(chars_per_line_estimate, 1))
    
    # Вычисляем необходимую высоту
    # Межстрочный интервал примерно 1.3 * размер_шрифта в пунктах
    line_height_pt = initial_size * 1.3
    line_height_inches = line_height_pt / 72.0  # Конвертируем в дюймы
    needed_height = estimated_lines * line_height_inches
    
    # Если текст не помещается, уменьшаем размер шрифта
    if needed_height > max_height:
        scale_factor = max_height / needed_height
        new_size = int(initial_size * scale_factor * 0.85)  # 0.85 для запаса
        return max(new_size, 10)  # Минимум 10pt для читаемости
    
    return initial_size

def add_text_to_shape(shape, text, font_size, bold=False, color=None, is_code=False, auto_fit=True):
    """Добавляет текст в shape с форматированием"""
    text_frame = shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Вычисляем оптимальный размер шрифта, если нужно
    if auto_fit and not is_code:
        # Получаем размеры shape
        max_width = shape.width
        max_height = shape.height
        font_size = calculate_font_size(text, max_width, max_height, font_size, is_code)
    
    # Если включен auto_fit, используем авторазмер (но только как запасной вариант)
    if auto_fit and not is_code:
        # Используем авторазмер для подгонки, но с ограничениями
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        text_frame.auto_size = None
    
    # Обработка списков и форматирования
    lines = text.split('\n')
    first_para = True
    
    for line in lines:
        if not line.strip() and not first_para:
            # Пустая строка - добавляем пустой параграф
            p = text_frame.add_paragraph()
            p.alignment = PP_ALIGN.LEFT
            continue
        
        if first_para:
            p = text_frame.paragraphs[0]
            first_para = False
        else:
            p = text_frame.add_paragraph()
        
        p.alignment = PP_ALIGN.LEFT
        
        # Определяем уровень списка
        list_level = 0
        if line.strip().startswith('- ') or line.strip().startswith('• '):
            line = line.replace('- ', '', 1).replace('• ', '', 1)
            list_level = 0
        elif line.strip().startswith('  - ') or line.strip().startswith('  • '):
            line = line.replace('  - ', '', 1).replace('  • ', '', 1)
            list_level = 1
        elif line.strip().startswith('    - ') or line.strip().startswith('    • '):
            line = line.replace('    - ', '', 1).replace('    • ', '', 1)
            list_level = 2
        
        p.level = list_level
        
        # Обработка жирного текста **text**
        parts = re.split(r'(\*\*[^*]+\*\*)', line)
        for part in parts:
            if part.startswith('**') and part.endswith('**'):
                run = p.add_run()
                run.text = part[2:-2]
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.name = 'Courier New' if is_code else FONT_NAME
                # Жирный текст делаем цветным
                run.font.color.rgb = ACCENT_COLOR if not is_code else (color if color else TEXT_COLOR)
            elif part.strip():
                run = p.add_run()
                run.text = part
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.name = 'Courier New' if is_code else FONT_NAME
                if color:
                    run.font.color.rgb = color
                else:
                    run.font.color.rgb = TEXT_COLOR

def add_gradient_background(slide, prs):
    """Добавляет градиентный фон к слайду"""
    # Создаем прямоугольник на весь слайд для градиента
    left = Inches(0)
    top = Inches(0)
    width = prs.slide_width
    height = prs.slide_height
    
    # Добавляем прямоугольник
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    
    # Устанавливаем градиент
    fill = shape.fill
    fill.gradient()
    fill.gradient_angle = 135.0  # Диагональный градиент
    
    # Настраиваем градиентные остановки
    # В python-pptx градиент обычно создается с двумя остановками
    gradient_stops = fill.gradient_stops
    
    # Первая остановка (начало - синий)
    if len(gradient_stops) > 0:
        gradient_stops[0].color.rgb = GRADIENT_START
        gradient_stops[0].position = 0.0
    
    # Вторая остановка (конец - фиолетовый)
    if len(gradient_stops) > 1:
        gradient_stops[1].color.rgb = GRADIENT_END
        gradient_stops[1].position = 1.0
    elif len(gradient_stops) == 1:
        # Если только одна остановка, используем средний цвет для визуального эффекта
        # Создаем второй прямоугольник с другим цветом для имитации градиента
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = GRADIENT_END
        shape2.line.fill.background()
        # Делаем полупрозрачным
        shape2.fill.transparency = 0.5
        # Отправляем на задний план
        slide.shapes._spTree.remove(shape2._element)
        slide.shapes._spTree.insert(2, shape2._element)
    
    # Убираем контур
    shape.line.fill.background()
    
    # Отправляем на задний план
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

def create_slide(prs, slide_data):
    """Создает слайд из данных"""
    # Используем пустой макет, чтобы избежать placeholders
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Добавляем градиентный фон
    add_gradient_background(slide, prs)
    
    # Добавляем белый фон для контента с тенью
    content_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.3), Inches(0.3),
        Inches(9.4), Inches(6.9)
    )
    content_bg.fill.solid()
    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    # Добавляем тонкую рамку для красоты
    content_bg.line.color.rgb = RGBColor(230, 230, 240)
    content_bg.line.width = Pt(1)
    
    # Добавляем декоративные элементы - маленькие круги в углах
    for corner_x, corner_y in [(0.2, 0.2), (9.8, 0.2), (0.2, 7.3), (9.8, 7.3)]:
        decor = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(corner_x), Inches(corner_y),
            Inches(0.15), Inches(0.15)
        )
        decor.fill.solid()
        decor.fill.fore_color.rgb = ACCENT_COLOR
        decor.line.fill.background()
        # Отправляем на задний план
        slide.shapes._spTree.remove(decor._element)
        slide.shapes._spTree.insert(-1, decor._element)
    
    # Заголовок слайда
    if slide_data['title']:
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1.2)
        title_box = slide.shapes.add_textbox(left, top, width, height)
        title_frame = title_box.text_frame
        title_frame.clear()
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.add_run()
        title_run.text = slide_data['title']
        title_run.font.size = Pt(TITLE_SIZE)
        title_run.font.bold = True
        title_run.font.name = FONT_NAME
        title_run.font.color.rgb = TITLE_COLOR
    
    # Подзаголовок (если есть)
    subtitle_top = Inches(1.8) if slide_data['title'] else Inches(0.8)
    if slide_data['subtitle']:
        left = Inches(0.5)
        top = subtitle_top
        width = Inches(9)
        height = Inches(0.8)
        subtitle_box = slide.shapes.add_textbox(left, top, width, height)
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_para.add_run()
        subtitle_run.text = slide_data['subtitle']
        subtitle_run.font.size = Pt(SUBTITLE_SIZE)
        subtitle_run.font.bold = True
        subtitle_run.font.name = FONT_NAME
        subtitle_run.font.color.rgb = SUBTITLE_COLOR
    
    # Содержимое слайда
    if slide_data['content']:
        # Вычисляем позицию и размер для контента
        content_top = Inches(2.8) if slide_data['subtitle'] else (Inches(2.0) if slide_data['title'] else Inches(1.0))
        content_bottom = Inches(6.7)  # Почти до конца слайда, оставляем небольшой отступ
        left = Inches(0.8)
        top = content_top
        width = Inches(8.4)
        height = content_bottom - content_top  # Динамическая высота
        
        content_box = slide.shapes.add_textbox(left, top, width, height)
        
        # Добавляем весь контент
        all_content_parts = []
        for content, is_code in slide_data['content']:
            all_content_parts.append((content, is_code))
        
        # Объединяем контент
        combined_lines = []
        for content, is_code in all_content_parts:
            combined_lines.append(content)
        
        combined_text = '\n\n'.join(combined_lines)
        # Используем моноширинный шрифт для блоков кода
        has_code = any(c[1] for c in all_content_parts)
        
        # Для блоков кода добавляем фон
        if has_code:
            # Добавляем фон для кода
            code_bg = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left - Inches(0.1), top - Inches(0.1),
                width + Inches(0.2), height + Inches(0.2)
            )
            code_bg.fill.solid()
            code_bg.fill.fore_color.rgb = CODE_BG_COLOR
            code_bg.line.fill.background()
            # Отправляем на задний план
            slide.shapes._spTree.remove(code_bg._element)
            slide.shapes._spTree.insert(-2, code_bg._element)
        
        add_text_to_shape(content_box, combined_text, CODE_SIZE if has_code else BODY_SIZE, 
                         is_code=has_code, color=TEXT_COLOR)
    
    return slide

def generate_pptx_from_markdown(md_file, output_file):
    """Генерирует PPTX презентацию из markdown файла"""
    print(f"📄 Чтение {md_file}...")
    try:
        slides = parse_markdown(md_file)
        print(f"   Найдено слайдов: {len(slides)}")
    except Exception as e:
        print(f"❌ Ошибка при парсинге markdown: {e}")
        import traceback
        traceback.print_exc()
        return
    
    # Создаем презентацию
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        print("📝 Генерация презентации...")
        
        for i, slide_data in enumerate(slides, 1):
            print(f"  Слайд {i}/{len(slides)}: {slide_data['title'] or 'Без заголовка'}")
            try:
                create_slide(prs, slide_data)
            except Exception as e:
                print(f"  ❌ Ошибка при создании слайда {i}: {e}")
                import traceback
                traceback.print_exc()
                continue
        
        # Сохраняем презентацию
        prs.save(output_file)
        print(f"\n✅ Презентация создана: {output_file}")
        print(f"   Размер: {os.path.getsize(output_file) / 1024:.1f} KB")
        print(f"   Количество слайдов: {len(slides)}")
    except Exception as e:
        print(f"❌ Ошибка при создании презентации: {e}")
        import traceback
        traceback.print_exc()


if __name__ == '__main__':
    parser = argparse.ArgumentParser(
        description='Генерация презентации проекта home-sentinel из presentation.md'
    )
    parser.add_argument(
        '--publish',
        action='store_true',
        help='Подготовить HTML для публикации в GitHub Pages (только для HTML генераторов)'
    )
    parser.add_argument(
        '--release',
        action='store_true',
        help='Создать GitHub Release с сгенерированной презентацией (DOCX/PPTX)'
    )
    parser.add_argument(
        '--version',
        type=str,
        help='Версия для GitHub Release (по умолчанию: timestamp)'
    )
    parser.add_argument(
        '--input',
        type=str,
        help='Путь к исходному markdown файлу (по умолчанию: docs/presentation.md)'
    )
    parser.add_argument(
        '--output',
        type=str,
        help='Путь к выходному файлу (по умолчанию: используется BUILD_DIR и шаблон из .env)'
    )
    
    args = parser.parse_args()
    
    # Загружаем конфигурацию
    env_vars = load_env_file()
    
    # Определяем пути
    project_root = get_project_root()
    
    # Путь к исходному файлу
    if args.input:
        md_file = args.input
    else:
        md_file = project_root / env_vars.get('DOCS_PRESENTATION_MD', 'docs/presentation.md')
    
    if not os.path.exists(md_file):
        print(f"❌ Файл не найден: {md_file}")
        sys.exit(1)
    
    # Путь к выходному файлу
    if args.output:
        output_file = args.output
    else:
        # Генерируем имя файла по шаблону
        filename, ext = generate_output_filename(str(md_file), env_vars=env_vars, output_ext='.pptx')
        build_dir = get_build_dir(env_vars)
        output_file = build_dir / f"{filename}{ext}"
    
    # Создаем директорию для выходного файла если нужно
    os.makedirs(os.path.dirname(output_file), exist_ok=True)
    
    # Генерируем презентацию
    generate_pptx_from_markdown(str(md_file), str(output_file))
    
    # Создаем GitHub Release если нужно (для DOCX/PPTX)
    if args.release:
        create_github_release([str(output_file)], env_vars, args.version)

